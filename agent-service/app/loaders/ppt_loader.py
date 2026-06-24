"""PowerPoint 加载器 — 提取每页标题、正文、备注"""

import logging

logger = logging.getLogger(__name__)


class PowerPointLoader:
    """读取 .pptx 所有页面内容，按 Slide 组织输出"""

    @staticmethod
    def supports(extension: str) -> bool:
        return extension.lower() in ("pptx", "ppt")

    def load(self, file_path: str) -> str:
        from pptx import Presentation

        logger.info("[PowerPointLoader] 开始解析: %s", file_path)

        prs = Presentation(file_path)
        parts = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_texts = [f"Slide {slide_idx}"]

            title = self._get_title(slide)
            if title:
                slide_texts.append(f"标题: {title}")

            body_texts = self._get_body_texts(slide)
            if body_texts:
                slide_texts.append("正文:")
                slide_texts.extend(f"  • {t}" for t in body_texts)

            notes = self._get_notes(slide)
            if notes and notes.strip():
                slide_texts.append(f"备注: {notes.strip()}")

            parts.append("\n".join(slide_texts))

        text = "\n\n".join(parts)
        logger.info("[PowerPointLoader] 解析完成, 页数=%d, 长度=%d", len(prs.slides), len(text))
        return text

    def _get_title(self, slide) -> str | None:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        return None

    def _get_body_texts(self, slide) -> list[str]:
        texts = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        texts.append(txt)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        return texts

    def _get_notes(self, slide) -> str | None:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            if tf:
                return tf.text.strip()
        return None